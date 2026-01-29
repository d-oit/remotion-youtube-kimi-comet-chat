#!/usr/bin/env python3
"""
PPTX Content Extractor

Extract text, images, and structure from PowerPoint (.pptx) files.
Outputs JSON for use in videos, documentation, or analysis.

Usage:
    python extract_pptx.py <input.pptx> [--output <output.json>]

Dependencies:
    pip install python-pptx
"""

import json
import sys
import os
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed.")
    print("Install with: pip install python-pptx")
    sys.exit(1)


def extract_text_from_element(element) -> dict:
    """Extract text and styling from a paragraph or run."""
    result = {"text": "", "style": "normal"}

    if hasattr(element, "text"):
        result["text"] = element.text.strip()

    # Check for bold/italic formatting
    if hasattr(element, "runs") and element.runs:
        for run in element.runs:
            if hasattr(run, "font"):
                if run.font.bold:
                    result["style"] = "bold"
                elif run.font.italic:
                    result["style"] = "italic"

    return result


def get_paragraph_style(paragraph) -> str:
    """Determine paragraph style based on indentation and level."""
    level = getattr(paragraph, "level", 0) or 0

    if level == 0:
        return "title"
    elif level == 1:
        return "subtitle"
    else:
        return "normal"


def extract_slide_content(slide) -> dict:
    """Extract all content from a single slide."""
    slide_data = {
        "slide_number": 0,
        "title": "",
        "layout": "content",
        "content": [],
        "notes": "",
        "images": [],
    }

    # Extract title
    if slide.shapes.title:
        slide_data["title"] = slide.shapes.title.text.strip()

    # Determine layout based on placeholders
    if slide.shapes.title and len(slide.placeholders) == 1:
        slide_data["layout"] = "title_only"
    elif slide.shapes.title and len(slide.placeholders) >= 2:
        if slide.placeholders[1].placeholder_format.type == 1:  # Body
            slide_data["layout"] = "title_content"

    # Extract body content
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE.SHAPE_TYPE_GROUP:
            continue

        # Skip title (already extracted)
        if shape == slide.shapes.title:
            continue

        # Skip notes page
        if shape.name and "Notes" in shape.name:
            continue

        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            if text and text != slide_data["title"]:
                # Determine if it's a bullet
                is_bullet = False
                if hasattr(shape, "paragraphs"):
                    for para in shape.paragraphs:
                        if hasattr(para, "level"):
                            if para.level > 0 or (
                                para.text.startswith("·")
                                or para.text.startswith("•")
                                or para.text.startswith("-")
                            ):
                                is_bullet = True
                                break

                content_item = {
                    "type": "bullet" if is_bullet else "paragraph",
                    "text": text,
                    "style": get_paragraph_style(shape)
                    if hasattr(shape, "paragraphs")
                    else "normal",
                }
                slide_data["content"].append(content_item)

        # Extract images
        if shape.shape_type == 13:  # MSO_SHAPE.PICTURE
            image_info = {
                "path": "",
                "slide": slide_data["slide_number"],
                "description": shape.name or "image",
            }
            if hasattr(shape, "image"):
                image_info["path"] = (
                    f"media/{shape.image.filename}"
                    if hasattr(shape.image, "filename")
                    else "image"
                )
                image_info["format"] = (
                    str(shape.image.format)
                    if hasattr(shape.image, "format")
                    else "unknown"
                )
            slide_data["images"].append(image_info)

    # Extract speaker notes
    if slide.has_notes_slide and slide.notes_slide:
        notes_text = (
            slide.notes_slide.notes_text_frame.text
            if slide.notes_slide.notes_text_frame
            else ""
        )
        slide_data["notes"] = notes_text.strip()

    return slide_data


def extract_presentation_metadata(prs) -> dict:
    """Extract presentation-level metadata."""
    metadata = {
        "title": "",
        "author": "",
        "slides_count": len(prs.slides),
        "created": "",
        "modified": "",
        "comments": "",
    }

    # Try to get properties
    if hasattr(prs, "core_properties"):
        props = prs.core_properties
        metadata["title"] = props.title or ""
        metadata["author"] = props.last_modified_by or props.creator or ""
        metadata["created"] = str(props.created)[:10] if props.created else ""
        metadata["modified"] = str(props.modified)[:10] if props.modified else ""

    return metadata


def extract_pptx(pptx_path: str, output_path: Optional[str] = None) -> dict:
    """
    Main extraction function.

    Args:
        pptx_path: Path to input .pptx file
        output_path: Optional path to save JSON output

    Returns:
        Dictionary with extracted presentation data
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"File not found: {pptx_path}")

    prs = Presentation(pptx_path)

    result = {
        "metadata": extract_presentation_metadata(prs),
        "slides": [],
        "images": [],
        "statistics": {"total_slides": 0, "total_images": 0, "total_notes": 0},
    }

    # Extract each slide
    for idx, slide in enumerate(prs.slides, start=1):
        slide_data = extract_slide_content(slide)
        slide_data["slide_number"] = idx
        result["slides"].append(slide_data)
        result["images"].extend(slide_data["images"])

        if slide_data["notes"]:
            result["statistics"]["total_notes"] += 1

    result["statistics"]["total_slides"] = len(result["slides"])
    result["statistics"]["total_images"] = len(result["images"])

    # Add segment mapping based on slide content
    result["segments"] = identify_segments(result["slides"])

    # Save output if specified
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        print(f"Extracted {result['metadata']['slides_count']} slides to {output_path}")

    return result


def identify_segments(slides: list) -> list:
    """Identify logical segments based on slide content."""
    segments = []
    current_segment = None

    for slide in slides:
        title = slide.get("title", "").lower()

        # Detect segment boundaries
        if "phase" in title or "step" in title or "workflow" in title:
            if current_segment:
                segments.append(current_segment)
            current_segment = {
                "name": slide["title"],
                "start_slide": slide["slide_number"],
                "end_slide": slide["slide_number"],
                "content": [slide["title"]]
                + [c["text"] for c in slide.get("content", [])],
            }
        elif current_segment:
            current_segment["end_slide"] = slide["slide_number"]
            current_segment["content"].append(slide["title"])

    # Don't forget last segment
    if current_segment:
        segments.append(current_segment)

    return segments


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <input.pptx> [--output <output.json>]")
        print("\nExample:")
        print("  python extract_pptx.py presentation.pptx")
        print("  python extract_pptx.py presentation.pptx --output storyboard.json")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = None

    if "--output" in sys.argv:
        idx = sys.argv.index("--output")
        if idx + 1 < len(sys.argv):
            output_path = sys.argv[idx + 1]

    try:
        result = extract_pptx(pptx_path, output_path)

        # Print summary
        print(f"\n=== Presentation Summary ===")
        print(f"Title: {result['metadata']['title']}")
        print(f"Slides: {result['metadata']['slides_count']}")
        print(f"Images: {result['statistics']['total_images']}")
        print(f"Notes: {result['statistics']['total_notes']}")

        if result.get("segments"):
            print(f"\n=== Identified Segments ===")
            for seg in result["segments"]:
                print(
                    f"  {seg['name']}: slides {seg['start_slide']}-{seg['end_slide']}"
                )

    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
